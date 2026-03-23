"""PowerPoint extraction using python-pptx."""

from pathlib import Path
from typing import Any

from app.core.logging import get_logger
from app.extraction.base import BaseExtractor, ExtractionResult, PageContent

logger = get_logger(__name__)


class PPTExtractor(BaseExtractor):
    SUPPORTED_EXTENSIONS = [".pptx", ".ppt"]

    async def extract(self, file_path: Path) -> ExtractionResult:
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            pages: list[PageContent] = []
            metadata: dict[str, Any] = {
                "total_slides": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
            }

            for idx, slide in enumerate(prs.slides, start=1):
                texts: list[str] = []
                slide_meta: dict[str, Any] = {
                    "slide_number": idx,
                    "layout": slide.slide_layout.name if slide.slide_layout else "",
                }

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)

                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            table_data.append([cell.text for cell in row.cells])
                        slide_meta.setdefault("tables", []).append(table_data)

                # Extract speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_meta["speaker_notes"] = notes

                content = "\n".join(texts)
                pages.append(PageContent(page_number=idx, content=content, metadata=slide_meta))

            return ExtractionResult(pages=pages, metadata=metadata)

        except Exception as e:
            logger.error("PowerPoint extraction failed for %s: %s", file_path, e)
            return ExtractionResult(success=False, error=str(e))
